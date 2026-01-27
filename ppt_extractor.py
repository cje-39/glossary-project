"""
PPT 병렬 텍스트 추출기
한국어/일본어 PPTX 파일에서 텍스트를 추출하여 엑셀로 병렬 정리

사용법:
    python ppt_extractor.py korean.pptx japanese.pptx

출력:
    ppt_parallel_text.xlsx (한국어, 일본어 병렬 텍스트)
"""

import sys
from pptx import Presentation
import pandas as pd
from pathlib import Path

def extract_text_from_ppt(ppt_path):
    """
    PPTX 파일에서 슬라이드별 텍스트 추출
    
    Returns:
        list: [{'slide_num': 1, 'texts': ['텍스트1', '텍스트2', ...]}, ...]
    """
    prs = Presentation(ppt_path)
    slides_data = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        
        # 슬라이드의 모든 shape를 순회
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            
            # 테이블 안의 텍스트도 추출
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text.strip())
        
        slides_data.append({
            'slide_num': slide_num,
            'texts': texts
        })
    
    return slides_data

def create_parallel_excel(korean_data, japanese_data, output_path='ppt_parallel_text.xlsx'):
    """
    한국어/일본어 데이터를 병렬로 정리하여 엑셀 생성
    """
    parallel_rows = []
    
    # 슬라이드 수가 다를 수 있으므로 최대값 사용
    max_slides = max(len(korean_data), len(japanese_data))
    
    for i in range(max_slides):
        ko_slide = korean_data[i] if i < len(korean_data) else {'slide_num': i+1, 'texts': []}
        ja_slide = japanese_data[i] if i < len(japanese_data) else {'slide_num': i+1, 'texts': []}
        
        # 텍스트 수가 다를 수 있으므로 최대값 사용
        max_texts = max(len(ko_slide['texts']), len(ja_slide['texts']))
        
        if max_texts == 0:
            # 빈 슬라이드도 기록
            parallel_rows.append({
                '슬라이드 번호': ko_slide['slide_num'],
                '텍스트 번호': 1,
                '한국어': '',
                '일본어': ''
            })
        else:
            for j in range(max_texts):
                ko_text = ko_slide['texts'][j] if j < len(ko_slide['texts']) else ''
                ja_text = ja_slide['texts'][j] if j < len(ja_slide['texts']) else ''
                
                parallel_rows.append({
                    '슬라이드 번호': ko_slide['slide_num'],
                    '텍스트 번호': j + 1,
                    '한국어': ko_text,
                    '일본어': ja_text
                })
    
    # DataFrame 생성 및 엑셀 저장
    df = pd.DataFrame(parallel_rows)
    
    # 엑셀 writer로 스타일 적용
    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='병렬텍스트')
        
        # 워크시트 가져오기
        worksheet = writer.sheets['병렬텍스트']
        
        # 열 너비 자동 조정
        worksheet.column_dimensions['A'].width = 12
        worksheet.column_dimensions['B'].width = 12
        worksheet.column_dimensions['C'].width = 50
        worksheet.column_dimensions['D'].width = 50
    
    print(f"✅ 엑셀 파일 생성 완료: {output_path}")
    print(f"   총 {len(parallel_rows)}개의 텍스트 쌍이 추출되었습니다.")
    
    return df

def main():
    if len(sys.argv) != 3:
        print("사용법: python ppt_extractor.py <한국어.pptx> <일본어.pptx>")
        print("예시: python ppt_extractor.py korean.pptx japanese.pptx")
        sys.exit(1)
    
    korean_path = sys.argv[1]
    japanese_path = sys.argv[2]
    
    # 파일 존재 확인
    if not Path(korean_path).exists():
        print(f"❌ 파일을 찾을 수 없습니다: {korean_path}")
        sys.exit(1)
    
    if not Path(japanese_path).exists():
        print(f"❌ 파일을 찾을 수 없습니다: {japanese_path}")
        sys.exit(1)
    
    print(f"📄 한국어 PPT 읽는 중: {korean_path}")
    korean_data = extract_text_from_ppt(korean_path)
    print(f"   → {len(korean_data)}개 슬라이드 발견")
    
    print(f"📄 일본어 PPT 읽는 중: {japanese_path}")
    japanese_data = extract_text_from_ppt(japanese_path)
    print(f"   → {len(japanese_data)}개 슬라이드 발견")
    
    print(f"\n📊 병렬 엑셀 생성 중...")
    df = create_parallel_excel(korean_data, japanese_data)
    
    print("\n✨ 완료!")

if __name__ == "__main__":
    main()
